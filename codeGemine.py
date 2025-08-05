import re
import streamlit as st
import tempfile
import google.generativeai as genai
from google.generativeai.types import GenerationConfig
from datetime import date
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION
import json
from googletrans import Translator
from huggingface_hub import InferenceClient
import pdfplumber
import os


def generate_image_hf(prompt, output_path):
    client = InferenceClient(
        provider="hf-inference",
        api_key=st.secrets["HF_API_KEY"]

    )

    # This returns a PIL.Image object
    image = client.text_to_image(
        prompt,
        model="stabilityai/stable-diffusion-3-medium-diffusers",
    )

    # Save PIL image to file
    image.save(output_path)
    return output_path

genai.configure(api_key=st.secrets["GOOGLE_API_KEY"])

def read_pdf(file_path):

    text = ''
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + '\n'
    return text.strip()

def read_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == '.docx':
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])

    elif ext == '.pdf':
        return read_pdf(file_path)
    else:
        raise ValueError("Unsupported file format: Only .docx and .pdf are supported.")


def build_prompt(text, slide_count = 6, include_visuals = False):
    remaining = slide_count - 3

    if include_visuals:
        main = remaining // 2
        extra = remaining % 2  # will be 1 if odd
        if extra == 1:
            main += 1
        else:
            slide_count = main + 3
        note = (
            "• Qeyd: Slayd sayı tam bölünmədiyi üçün sonuncu Əsas slayd `visual.type = 'none'` olaraq təyin edilməlidir."
            if extra == 1 else ""
        )
    else:
        main = remaining
        note = ""
    return f"""
Sənə bir sənədin mətni təqdim olunur. Bu mətni təhlil et və təqdimat üçün aşağıdakı struktura uyğun slayd formatında hazırla:

TƏQDİMATIN STRUKTURU:
1. Başlıq – Təqdimatın adı 
2. Giriş – Məqsəd və təqdim olunacaq mövzu üzrə qısa xülasə
3. Əsas Slaydlar – Mətndəki əsas mövzulardan hər biri üçün ayrıca slayd:
   • Hər biri unikal başlığa sahib olmalıdır  
   • Hər slaydda **4 əsas bənd** olmalıdır  
   • Slayd sayı: İstifadəçi {slide_count} slayd istəmişdir, bu saydan **1 başlıq**, **1 giriş**, **1 tövsiyə** slayd çıxıldıqdan sonra qalan **{main}** slayd əsas mövzular üçün istifadə olunmalıdır. 
4. Tövsiyə – Gələcək inkişaf və təkmilləşdirmə üçün 4–5 bəndlik təkliflər.

QAYDALAR:
- İstifadəçi {slide_count} slayd istəmişdir. Mətni analiz et və bu sayda slayda uyğun şəkildə ayır.
- Hər slayd üçün JSON obyektində aşağıdakı sahələri yaz:
  - type: "title" | "intro" | "main" | "recommendation"
  - Başlıq üçün (type = "title"):
    • title (təqdimatın adı)
  - Giriş üçün (type = "intro"):
    • aim (təqdimatın məqsədi)
    • summary (layihənin qısa xülasəsi, 3-4 cümlə)
  - Əsas slaydlar üçün (type = "main"):
    • title (slaydın başlığı)
    • point1, point2, point3, point4 (hər biri əsas məzmun bəndləri)
    • visual (vizual təklif üçün JSON obyekti, aşağıdakı formatda)
  - Tövsiyə slaydı üçün (type = "recommendation"):
    • recommendation1, recommendation2, recommendation3, recommendation4, (optional) recommendation5
- Hər slayd üçün:
  • Mətni slayd sayına uyğun şəkildə bərabər böl.
  • Cümlə dəyərlərinin içində qaçırılmamış qoşa dırnaq işarələrindən istifadə etmə. Daxili dırnaq işarələri üçün tək dırnaqdan istifadə et.
  • Lazımsız təkrarlardan, çox uzun cümlələrdən qaç.
  • Slayd dili **rəsmi və aydın olmalıdır**.
  • Əgər statistik və ya əsas nəticələr varsa, Əsas Göstəricilər slaydına daxil et.
  • Slaydların məzmunu yalnız sənədə əsaslanmalıdır, əlavə məlumat əlavə etmə.
  {note}
 
- Vizual təklif JSON formatı (mümkün olduqca fərqli növ vizuallar əlavə et, type image daxil olmaqla.):
    ```json
    {{
        "type": "none" | "image" | "bar" | "pie" | "line",
        "title": "Vizualın başlığı",
        "description": "Əgər type 'image'dirsə, burada şəkilin ətraflı təsviri verilir. Digər hallarda boş saxlanılır.",
        "xlabel": "X oxunun etiketi (əgər tətbiq olunursa)",
        "ylabel": "Y oxunun etiketi (əgər tətbiq olunursa)",
        "x": ["X oxundakı dəyərlər (əgər varsa)"],
        "y": ["Y oxundakı dəyərlər (əgər varsa)"],
        "labels": ["Dilim adları (əgər 'pie' tipi tətbiq olunursa)"],
        "sizes": ["Dilim ölçüləri (əgər 'pie' tipi tətbiq olunursa)"]
    }}
    ```

    - `type` sahəsi yalnız aşağıdakı dəyərlərdən biri ola bilər: "none", "image", "bar", "pie", "line"
    - `type` = "image" olduqda, `description` sahəsi vacibdir və şəkilin məzmununu izah etməlidir.
    - `type` = "bar" və ya "line" olduqda x, y, xlabel, ylabel sahələri dolu olmalıdır.
    -  Əgər uyğun vizual yoxdursa, `type` dəyəri `"none"` olmalıdır.
    - `type` = "pie" olduqda labels və sizes sahələri dolu olmalı, `x`, `y`, `xlabel`, `ylabel` isə boş buraxılmalıdır.
    - `type` = "none" olduqda digər sahələr boş buraxılmalıdır.

TƏHLİL EDİLƏCƏK SƏNƏDİN MƏTNI:
\"\"\"
{text}
\"\"\"

CAVABI BU FORMATA JSON ARRAY KİMİ QAYTAR:

```json
[
  {{
    "type": "title",
    "title": "Təqdimatın adı",
  }},
  {{
    "type": "intro",
    "aim": "Təqdimatın məqsədi",
    "summary": "Layihənin qısa xülasəsi"
  }},
  {{
    "type": "main",
    "title": "Əsas mövzunun başlığı",
    "point1": "...",
    "point2": "...",
    "point3": "...",
    "point4": "...",
    "visual": {{
      "type": "bar",
      "title": "",
      "description": "",
      "xlabel": "",
      "ylabel": "",
      "x": [],
      "y": [],
      "labels": [],
      "sizes": []
    }}
  }},
  ...
  {{
    "type": "recommendation",
    "recommendation1": "...",
    "recommendation2": "...",
    "recommendation3": "...",
    "recommendation4": "...",
    "recommendation5": "..."
  }}
]
```
"""


def get_presentation(text, slide_count=6, model_name= 'gemini-1.5-flash', include_visuals = False):
    prompt = build_prompt(text, slide_count, include_visuals)
    print(prompt)
    model = genai.GenerativeModel(
        model_name=model_name,
        system_instruction="Sən təqdimat üzrə Azərbaycan dilində AI asistentsən."
    )

    generation_config = GenerationConfig(
        temperature=0.3,  # Controls randomness. Lower values are more deterministic.
    )

    try:
        # Send the prompt to the Gemini model
        response = model.generate_content(
            contents=[
                {"role": "user", "parts": [{"text": prompt}]}
            ],
            generation_config=generation_config
        )

        if response.candidates and response.candidates[0].content and response.candidates[0].content.parts:
            return response.candidates[0].content.parts[0].text
        else:
            print("Error: Model did not return expected content structure.")
            return "Content generation failed."

    except Exception as e:
        print(f"An error occurred during content generation: {e}")
        return f"Error: {e}"


def parse_gpt_response(response_text):
    try:

        match = re.search(r'\[.*\]', response_text, re.DOTALL)
        if match:
            cleaned_text = match.group(0).strip()

        slides = json.loads(cleaned_text)


    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON format: {e}")

    # Basic validation: check it's a list, and each slide has required fields depending on type
    if not isinstance(slides, list):
        raise ValueError("Expected a JSON array (list) of slides.")

    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            raise ValueError(f"Slide {i} is not a JSON object.")
        slide_type = slide.get("type")
        if slide_type not in {"title", "intro", "main", "recommendation"}:
            raise ValueError(f"Slide {i} has unknown type: {slide_type}")

        # Example minimal checks per slide type
        if slide_type == "title":
            if "title" not in slide:
                raise ValueError(f"Slide {i} of type 'title' missing required fields.")
        elif slide_type == "intro":
            if "aim" not in slide or "summary" not in slide:
                raise ValueError(f"Slide {i} of type 'intro' missing required fields.")
        elif slide_type == "main":
            required_fields = {"title", "point1", "point2", "point3", "point4", "visual"}
            if not required_fields.issubset(slide.keys()):
                missing = required_fields - set(slide.keys())
                raise ValueError(f"Slide {i} of type 'main' missing fields: {missing}")
            # Visual should be a dict
            if not isinstance(slide["visual"], dict):
                raise ValueError(f"Slide {i} 'visual' must be an object.")
        elif slide_type == "recommendation":
            # At least 4 recommendations required
            rec_fields = [f"recommendation{i}" for i in range(1, 6)]
            present_recs = [f for f in rec_fields if f in slide]
            if len(present_recs) < 4:
                raise ValueError(f"Slide {i} of type 'recommendation' requires at least 4 recommendations.")

    return slides




def add_title_slide(prs, slide):


    today = date.today()
    formatted_date = today.strftime("%d/%m/%Y")
    s = prs.slides[0]

    title_shape = s.shapes.title
    if title_shape and title_shape.text_frame:
        p = title_shape.text_frame.paragraphs[0]
        if p.runs:
            p.runs[0].text = slide['title']
        else:
            p.text = slide['title']

    for placeholder in s.placeholders:
        if "Tarix" in placeholder.text:
            text_frame = placeholder.text_frame
            p = text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text = f"Tarix: {formatted_date}"
            else:
                p.text = f"Tarix: {formatted_date}"
            break


def add_intro_slide(prs, slide):
    s = prs.slides[1] # Target the second slide
    vertical_margin = Inches(0.2)

    for shape in s.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            if text == "Layihənin məzmunu":

                if 'summary' in slide and slide['summary']:
                    new_text_box_width = shape.width
                    left = shape.left
                    top = shape.top + shape.height + vertical_margin

                    initial_height_for_autosize = Inches(3)

                    txBox = s.shapes.add_textbox(left, top, new_text_box_width, initial_height_for_autosize)
                    tf = txBox.text_frame

                    tf.word_wrap = True
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0

                    p_content = tf.paragraphs[0]
                    p_content.clear()
                    run_content = p_content.add_run()
                    run_content.text = slide['summary']
                    run_content.font.size = Pt(17)
                    run_content.font.bold = False

                else:
                    print("Summary data is missing or empty.")

            elif text == "Məqsəd":



                if 'aim' in slide and slide['aim']:
                    new_text_box_width = shape.width
                    left = shape.left
                    top = shape.top + shape.height + vertical_margin

                    initial_height_for_autosize = Inches(3)

                    txBox = s.shapes.add_textbox(left, top, new_text_box_width, initial_height_for_autosize)
                    tf = txBox.text_frame

                    tf.word_wrap = True
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0

                    p_content = tf.paragraphs[0]
                    p_content.clear()
                    run_content = p_content.add_run()
                    run_content.text = slide['aim']
                    run_content.font.size = Pt(17)
                    run_content.font.bold = False



                else:
                    print("Aim data is missing or empty.")
        else:
            print("    Does not have a text frame.")


def convert_sizes(sizes):
    result = []
    for s in sizes:
        s = s.strip()
        if s.endswith('%'):

            value = float(s[:-1])
        else:
            value = float(s)
        result.append(value)
    return result


def add_main_slide(prs, slide):
    layout = prs.slide_layouts[12]
    s = prs.slides.add_slide(layout)


    s.shapes.title.text = slide['title']
    current_point_idx = 1

    for shape in s.shapes:

        if shape.is_placeholder:
            current_text_in_shape = shape.text.strip()
            if current_text_in_shape == "" and current_point_idx <= 4:

                point = slide.get(f'point{current_point_idx}', '')

                if point:
                    tf = shape.text_frame
                    tf.clear()

                    p = tf.add_paragraph()
                    p.text = point

                    p.font.size = Pt(17)

                    print(f"Populated shape with 'xxxx' (Point {current_point_idx}): '{point}'")
                else:
                    shape.text_frame.clear()
                    print(f"No content for point{current_point_idx}, cleared 'xxxx' from shape.")

                current_point_idx += 1



    visual = slide.get('visual', {})
    if visual.get('type') and visual['type'] != 'none':

        visual_slide_layout = prs.slide_layouts[3]
        visual_s = prs.slides.add_slide(visual_slide_layout)
        visual_s.shapes.title.text = f"{slide['title']} - {visual.get('title', 'Visual')}"



        if visual["type"] in ["bar", "line"]:
            add_chart(visual_s, visual["type"], visual["title"],
                      visual["x"], list(map(float, visual["y"])),
                      visual.get("xlabel", ""), visual.get("ylabel", ""))

        elif visual["type"] == "pie":
            sizes = convert_sizes(visual["sizes"])
            add_chart(visual_s, visual["type"], visual["title"],
                      labels=visual['labels'], sizes=sizes)


        elif visual["type"] == "image":
            try:
                image_path = f"generated_image_{visual['title'][:10].replace(' ', '_')}.png"
                translator = Translator()
                translated = translator.translate(visual["description"], src='az', dest='en')
                english_description = translated.text
                print(english_description)
                generate_image_hf(english_description, image_path)

                try:
                    placeholder = visual_s.placeholders[1]
                    left = placeholder.left
                    top = placeholder.top
                    width = placeholder.width
                    height = placeholder.height

                    # Add the picture at placeholder's position and size
                    visual_s.shapes.add_picture(image_path, left, top, width=width, height=height)
                except IndexError:
                    print("Slide does not have a second placeholder.")
                    insert_text_or_fallback(visual_s, f"[Şəkil təsviri: {visual['description']}]")

            except Exception as e:
                print(f"Error generating image: {e}")
                insert_text_or_fallback(visual_s, f"[Şəkil təsviri: {visual['description']}]")
        else:
            insert_text_or_fallback(visual_s, f"[Vizual növü '{visual['type']}' hələ dəstəklənmir]")
            
def insert_text_or_fallback(slide, text):
    try:
        placeholder = slide.placeholders[1]  # second placeholder
        placeholder.text = text
    except IndexError:
        textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1))
        textbox.text_frame.text = text



def add_recommendation_slide(prs, slide):
    layout = prs.slide_layouts[3]
    s = prs.slides.add_slide(layout)
    title_shape = s.shapes.title


    if title_shape:
        title_shape.text = "Növbəti addımlar"

    target_tf = None

    for shape in s.shapes:

        if not shape.has_text_frame:
            continue

        paragraphs = shape.text_frame.paragraphs
        if not paragraphs:
            continue

        first_text = paragraphs[0].text.strip()

        if first_text == "":
            target_tf = shape.text_frame


    if target_tf:

        target_tf.clear()

        for i in range(1, 6):
            rec = slide.get(f'recommendation{i}', None)
            if rec:
                p = target_tf.add_paragraph()
                p.text = rec
                p.level = 1
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.font.size = Pt(21)

                print(f"Added recommendation{i}: {rec}")
            else:
                print(f"No data for recommendation{i}.")



def delete_slide(prs, slide_index):
    slide_id = prs.slides._sldIdLst[slide_index]
    prs.slides._sldIdLst.remove(slide_id)

def generate_pptx(slides, output_filename="presentation.pptx"):
    prs = Presentation("format_new.pptx")

    for slide in slides:
        t = slide.get('type')
        if t == 'title':
            add_title_slide(prs, slide)
        elif t == 'intro':
            add_intro_slide(prs, slide)
        elif t == 'main':
            add_main_slide(prs, slide)
        elif t == 'recommendation':
            add_recommendation_slide(prs, slide)


    delete_slide(prs, 3)
    delete_slide(prs, 2)



    prs.save(output_filename)
    print(f"Presentation saved as '{output_filename}'")



def add_chart(slide, chart_type, chart_title, x=None, y=None, xlabel=None, ylabel=None, labels=None, sizes=None):
    chart_data = CategoryChartData()

    if chart_type == "pie":
        if labels is not None and sizes is not None:
            chart_data.categories = labels
            chart_data.add_series("", sizes)
        else:
            print("Warning: Pie chart requires 'labels' and 'sizes' data.")
            return  # Exit
    elif chart_type in ["bar", "line"]:
        if x is not None and y is not None:
            chart_data.categories = x
            chart_data.add_series("", y)
        else:
            print(f"Warning: {chart_type} chart requires 'x' and 'y' data.")
            return
    else:
        print(f"Warning: Unsupported chart type: {chart_type}")
        return



    if chart_type == "bar":
        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
    elif chart_type == "line":
        chart_type_enum = XL_CHART_TYPE.LINE
    elif chart_type == "pie":
        chart_type_enum = XL_CHART_TYPE.PIE
    else:
        return


    placeholder = slide.placeholders[1]

    if placeholder:
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
    else:
        # fallback position
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(6)
        height = Inches(3)



    chart = slide.shapes.add_chart(chart_type_enum, left, top, width, height, chart_data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title
    chart.has_legend = False


    for series in chart.series:
        series.data_labels.show_value = True

    if chart_type == "pie":
            series.data_labels.show_category_name = True
            series.data_labels.show_percentage = True
            series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END


    if chart_type in ["bar", "line"]:
        category_axis = chart.category_axis
        value_axis = chart.value_axis

        if xlabel:
            category_axis.has_title = True
            category_axis.axis_title.text_frame.text = xlabel

        if ylabel:
            value_axis.has_title = True
            value_axis.axis_title.text_frame.text = ylabel






def streamlit():
    st.title("Sənəddən Təqdimat Yaratma")

    # File upload
    uploaded_file = st.file_uploader("PDF və ya DOCX faylını yükləyin", type=["pdf", "docx"])

    slide_count = st.number_input("Slaydların sayı", min_value=5, step=1)

    include_visuals = st.radio(
        "Vizuaları ümumi slayd sayına daxil edək?",
        ("Bəli", "Xeyr"),
        index=1,
        help="Vizual elementləri slayd sayına daxil etmək üçün 'Bəli' seçin."
    )

    if uploaded_file:
        generate_btn = st.button("PPTX Yarat")
    else:
        generate_btn = False

    if "pptx_bytes" not in st.session_state:
        st.session_state.pptx_bytes = None
    if "generation_done" not in st.session_state:
        st.session_state.generation_done = False

    if uploaded_file and generate_btn:
        try:
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[-1]) as tmp:
                tmp.write(uploaded_file.read())
                tmp_path = tmp.name

            # Step 1: Extract text
            with st.spinner("Fayl oxunur və təqdimat hazırlanır..."):
                doc_text = read_file(tmp_path)
                gpt_response = get_presentation(doc_text, slide_count, include_visuals=(include_visuals == "Bəli"))
                slides = parse_gpt_response(gpt_response)

                output_filename = "generated_presentation.pptx"
                generate_pptx(slides, output_filename)

                with open(output_filename, "rb") as f:
                    st.session_state.pptx_bytes = f.read()

            st.session_state.generation_done = True
            st.success("Təqdimat uğurla yaradıldı!")

        except Exception as e:
            st.error(f"Error: {e}")

    if st.session_state.generation_done and st.session_state.pptx_bytes:
        st.download_button(
            label="PPTX Faylını Yüklə",
            data=st.session_state.pptx_bytes,
            file_name="generated_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )


if __name__ == "__main__":
    streamlit()