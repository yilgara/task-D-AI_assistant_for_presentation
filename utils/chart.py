from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION





def add_chart(slide, chart_type, chart_title, x=None, y=None, xlabel=None, ylabel=None, labels=None, sizes=None):
    chart_data = CategoryChartData()

    if chart_type == "pie":
        if labels is not None and sizes is not None:
            chart_data.categories = labels
            chart_data.add_series("", sizes)
        else:
            print("Warning: Pie chart requires 'labels' and 'sizes' data.")
            return  # Exit
    elif chart_type in ["bar", "line"]:
        if x is not None and y is not None:
            chart_data.categories = x
            chart_data.add_series("", y)
        else:
            print(f"Warning: {chart_type} chart requires 'x' and 'y' data.")
            return
    else:
        print(f"Warning: Unsupported chart type: {chart_type}")
        return



    if chart_type == "bar":
        chart_type_enum = XL_CHART_TYPE.COLUMN_CLUSTERED
    elif chart_type == "line":
        chart_type_enum = XL_CHART_TYPE.LINE
    elif chart_type == "pie":
        chart_type_enum = XL_CHART_TYPE.PIE
    else:
        return


    placeholder = slide.placeholders[1]

    if placeholder:
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
    else:
        # fallback position
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(6)
        height = Inches(3)



    chart = slide.shapes.add_chart(chart_type_enum, left, top, width, height, chart_data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title
    chart.has_legend = False


    for series in chart.series:
        series.data_labels.show_value = True

    if chart_type == "pie":
            series.data_labels.show_category_name = True
            series.data_labels.show_percentage = True
            series.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END


    if chart_type in ["bar", "line"]:
        category_axis = chart.category_axis
        value_axis = chart.value_axis

        if xlabel:
            category_axis.has_title = True
            category_axis.axis_title.text_frame.text = xlabel

        if ylabel:
            value_axis.has_title = True
            value_axis.axis_title.text_frame.text = ylabel